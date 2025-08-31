"""
mcp_document_processing.py
FastMCP server exposing a Document‑processing tool
that works without the camel package.
"""

# --------------------------------------------------------------------------- #
#  Imports
# --------------------------------------------------------------------------- #
import asyncio, os, io, json, subprocess
from typing import Tuple, Optional, List, Literal

from loguru import logger
from retry import retry

from mcp.server.fastmcp import FastMCP
import anyio

# --- your own helper toolkits ------------------------------------------------ #
#   (provide these scripts in the Python path)
from image_tool import ask_question_about_image
from excel_tool import ExcelToolkit
from video_tool import ask_question_about_video
# --- third‑party libs already used ------------------------------------------ #
import assemblyai as aai
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from PIL import Image
from docx2markdown._docx_to_markdown import docx_to_markdown
from chunkr_ai import Chunkr
import xmltodict
import nest_asyncio
nest_asyncio.apply()

from dotenv import load_dotenv
load_dotenv(".env")


# --------------------------------------------------------------------------- #
#  Toolkit implementation (no camel.BaseToolkit!)
# --------------------------------------------------------------------------- #
class DocumentProcessingToolkit:
    """
    This tool exposes a **general‑purpose document‑processing endpoint** that
    converts almost any common file you point it to into **clean, readable
    text or Markdown**.  It is useful whenever an agent needs to “look inside”
    an arbitrary file before reasoning over its contents.

    Conceptually, you can think of it as an *all‑in‑one* “open the file and
    give me the text” utility:

    • **Images (.jpg / .jpeg / .png)** – runs a vision model and returns a
    detailed caption.  
    • **Audio (.mp3 / .wav / .m4a)** – performs automatic transcription.  
    • **PowerPoint (.pptx)** – pulls every textbox, captions each embedded
    image, and preserves the slide order.  
    • **Spreadsheets (.xls / .xlsx / .csv)** – dumps cell values in a
    readable, row‑wise layout.  
    • **ZIP archives** – unpacks the archive and lists the extracted files.  
    • **Plain text‑like formats (.py / .txt)** – simply reads the file.  
    • **JSON, JSONL, JSON‑LD** – returns the parsed JSON structure.  
    • **XML** – converts to a Python dict (falls back to raw XML on error).  
    • **Word (.docx)** – converts the entire document to Markdown.  
    • **Video (.mov)** – generates a descriptive narration of the clip.  
    • **PDF or any other format** – attempts Chunkr AI extraction first, and
    then a plain‑text PDF fallback if Chunkr fails.

    Typical downstream tasks include:

    - Letting an LLM **summarise** or **answer questions about** a
    presentation, spreadsheet, contract, or research paper.  
    - Quickly **indexing** large document batches for semantic search.  
    - **Captioning media assets** (images & video) to improve accessibility.  
    - Turning “opaque” binary files into human‑readable text for diffing or
    version control.

    ### Args
    `document_path` *(str)* – A fully‑qualified **local file path** that the
    server can access (e.g. `/home/user/input/report.pdf`).  Network URLs are
    *not* accepted.

    ### Returns
    `str` – On success, a **plain‑text or Markdown** representation of the
    file’s meaningful content.  
    If the file type is unsupported or an extraction error occurs, the tool
    raises an exception containing a diagnostic message.
    """

    def __init__(self, cache_dir: Optional[str] = None):
        self.excel_tool = ExcelToolkit()
        self.cache_dir = cache_dir or "tmp/"

    # --------------------------------------------------------------------- #
    #  Public façade
    # --------------------------------------------------------------------- #
    @retry(Exception,tries=5, delay=2, backoff=2)
    def extract_document_content(self, document_path: str) -> Tuple[bool, str]:
        logger.debug(f"[extract_document_content] {document_path=}")

        # 1. Images ----------------------------------------------------------------
        if document_path.lower().endswith((".jpg", ".jpeg", ".png")):
            caption = asyncio.run(
                ask_question_about_image(
                    document_path,
                    "Please make a detailed caption about the image."
                )
            )
            return True, caption

        # 2. Audio -----------------------------------------------------------------
        if document_path.lower().endswith((".mp3", ".wav", ".m4a")):
            aai.settings.api_key = os.getenv("ASSEMBLYAI_API_KEY")
            config = aai.TranscriptionConfig(speech_model=aai.SpeechModel.best)
            transcript = aai.Transcriber(config=config).transcribe(document_path)
            logger.info(transcript.text)
            if transcript.status == "error":
                raise RuntimeError(f"Transcription failed: {transcript.error}")
            return True, transcript.text

        # 3. PPTX ------------------------------------------------------------------
        if document_path.lower().endswith(".pptx"):
            return True, asyncio.run(self._extract_pptx(document_path))

        # 4. Spreadsheets -----------------------------------------------------------
        if document_path.lower().endswith((".xls", ".xlsx", ".csv")):
            return True, self.excel_tool.extract_excel_content(document_path)

        # 5. Zip --------------------------------------------------------------------
        if document_path.lower().endswith(".zip"):
            return True, f"The extracted files are: {self._unzip_file(document_path)}"

        # 6. Simple text‑like formats ----------------------------------------------
        simple_readers = {
            ".py":  lambda p: open(p, encoding="utf‑8").read(),
            ".txt": lambda p: open(p, encoding="utf‑8").read(),
        }
        if any(document_path.lower().endswith(ext) for ext in simple_readers):
            reader = simple_readers[os.path.splitext(document_path)[1]]
            return True, reader(document_path)

        # 7. JSON                                                                   #
        if document_path.lower().endswith((".json", ".jsonl", ".jsonld")):
            return True, self._extract_json(document_path, encoding="utf‑8")
        

        # 8. XML                                                                    #
        if document_path.lower().endswith(".xml"):
            data = open(document_path, encoding="utf‑8").read()
            try:
                return True, xmltodict.parse(data)
            except Exception:
                return True, data

        # 9. DOCX → markdown -------------------------------------------------------
        if document_path.lower().endswith(".docx"):
            md_path = f"{os.path.basename(document_path)}.md"
            docx_to_markdown(document_path, md_path)
            return True, open(md_path, encoding="utf‑8").read()

        # 10. MOV video ------------------------------------------------------------
        if document_path.lower().endswith(".mov"):
            description = ask_question_about_video(
                document_path, "Please make a detailed description about the video."
            )
            return True, description

        # 11. Fallback – Chunkr / PDF text -----------------------------------------
        return self._try_chunkr_then_fallback(document_path)

    # ------------------------------------------------------------------------- #
    #  helpers
    # ------------------------------------------------------------------------- #
    def _extract_json(self, json_path: str, encoding: str = "utf‑8") -> str:
        with open(json_path, 'r', encoding=encoding) as f:
            if json_path.lower().endswith((".json",".jsonld")):
                return json.load(f)  
            elif json_path.lower().endswith(".jsonl"):
                return [json.loads(line) for line in f]                    

    async def _extract_pptx(self, pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        base = pptx_path.rsplit(".", 1)[0]
        out = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            txt = [f"Page {slide_idx}"]
            captions = []
            img_count = 0

            for shape_idx, shape in enumerate(slide.shapes):
                if getattr(shape, "text", "").strip():
                    txt.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    img = Image.open(io.BytesIO(shape.image.blob))
                    img_path = f"{base}_slide_{slide_idx}_img_{shape_idx}.png"
                    img.save(img_path)
                    captions.append(
                        f"Image {img_count}: "
                        + await ask_question_about_image(
                            img_path, "Please make a detailed caption about the image."
                        )
                    )

            out.append("\n".join(txt + captions))

        return "\n\n".join(out)

    def _try_chunkr_then_fallback(self, path: str) -> Tuple[bool, str]:
        try:
            text = asyncio.run(
                self._extract_with_chunkr(path, output_format="markdown")
            )
            return True, text
        except Exception as e:
            logger.warning(f"Chunkr failed: {e}")
            if path.lower().endswith(".pdf"):
                try:
                    from PyPDF2 import PdfReader
                    text = "".join(
                        p.extract_text() for p in PdfReader(open(path, "rb")).pages
                    )
                    return True, text
                except Exception as e2:
                    return False, f"PDF fallback failed: {e2}"
            return False, f"Unsupported file type or processing error: {e}"

    async def _extract_with_chunkr(
        self, path: str, output_format: Literal["json", "markdown"] = "markdown"
    ) -> str:
        chunkr = Chunkr(api_key=os.getenv("CHUNKR_API_KEY"))
        result = await chunkr.upload(path)

        if result.status == "Failed":
            raise RuntimeError(result.message)

        out_path = f"{os.path.basename(path)}.{ 'json' if output_format=='json' else 'md' }"
        (result.json if output_format == "json" else result.markdown)(out_path)
        return open(out_path, encoding="utf‑8").read()

    def _unzip_file(self, zip_path: str) -> List[str]:
        dest = os.path.join(self.cache_dir, os.path.splitext(os.path.basename(zip_path))[0])
        os.makedirs(dest, exist_ok=True)
        subprocess.run(["unzip", "-o", zip_path, "-d", dest], check=True)
        return [os.path.join(r, f) for r, _, fs in os.walk(dest) for f in fs]


# --------------------------------------------------------------------------- #
#  FastMCP server
# --------------------------------------------------------------------------- #
mcp = FastMCP("document_processing")
toolkit = DocumentProcessingToolkit()


@mcp.tool()
async def process_document(document_path: str) -> str:
    """
    Process a document at the given *document_path*. The document can be multimedia
    (image, audio, video), a presentation (PPTX), a spreadsheet, a ZIP archive,
    a text file, JSON, XML, Word document, or PDF.
   Return the extracted text / markdown representation of *document_path*.
   supported formats include .png, .jpeg, .jpg, .mp3, .m4a, .PPTX, .xlsx, .csv, .txt, .json, .jsonl, .jsonld, .zip, .xml, .docx, .mov, and .pdf.
    """
    success, content = await anyio.to_thread.run_sync(
        toolkit.extract_document_content, document_path
    )
    if not success:
        raise ValueError(content)
    return content


# --------------------------------------------------------------------------- #
#  Entrypoint
# --------------------------------------------------------------------------- #
if __name__ == "__main__":
    mcp.run(transport="stdio")
